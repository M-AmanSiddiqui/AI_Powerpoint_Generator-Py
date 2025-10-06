from flask import Flask, render_template_string, request, send_file
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import re

# -----------------------------
# 🔑 API Key
# -----------------------------
from dotenv import load_dotenv
import os

load_dotenv()
api_key = os.getenv("API_KEY")

if not api_key:
    raise ValueError("❌ API key not found! Please set it in your .env file")

genai.configure(api_key=api_key)
model = genai.GenerativeModel("gemini-2.0-flash")
app = Flask(__name__)

# -----------------------------
# 🎨 Apply Background Color
# -----------------------------
def apply_custom_color(slide, color_hex):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    color_hex = color_hex.lstrip("#")
    r, g, b = tuple(int(color_hex[i:i + 2], 16) for i in (0, 2, 4))
    fill.fore_color.rgb = RGBColor(r, g, b)

# -----------------------------
# 🧠 Generate Slide Content
# -----------------------------
def get_generated_content(topic, slides, theme, content_type, count):
    try:
        if content_type == "bullets":
            prompt = f"""
            Generate exactly {slides} PowerPoint slides about "{topic}".
            Each slide must have a heading like "Slide 1:", "Slide 2:" etc.
            Each slide should contain:
            - Title: <title>
            - Exactly {count} short, natural bullet points (max 2 lines each)
            Use clear, human-like wording. Keep bullets concise.
            Style: {theme}.
            """
        else:
            prompt = f"""
            Generate exactly {slides} PowerPoint slides about "{topic}".
            Each slide must have a heading like "Slide 1:", "Slide 2:" etc.
            Each slide should contain:
            - Title: <title>
            - Exactly {count} short paragraphs (max 3 lines each)
            Do NOT use bullet points.
            Keep text concise and clear.
            Style: {theme}.
            """

        response = model.generate_content(prompt)
        text = getattr(response, "text", str(response))

        # Clean up formatting
        clean_text = re.sub(r"[*#>\-]+", "", text)
        clean_text = re.sub(r"(Bullet Points:|Visuals:|Content:|Points:)", "", clean_text, flags=re.IGNORECASE)
        clean_text = re.sub(r"\n{3,}", "\n\n", clean_text)
        clean_text = clean_text.strip()

        # ✅ Split slides safely
        slides_data = re.split(r"(?:Slide\s*\d+[:\-]?)", clean_text)
        slides_data = [s.strip() for s in slides_data if s.strip()]

        # If AI didn't label slides → auto split by bullet/paragraph count
        if len(slides_data) != slides:
            print("⚠️ Auto-splitting fallback triggered.")
            lines = [line.strip() for line in clean_text.split("\n") if line.strip()]
            chunk_size = max(1, len(lines) // slides)
            slides_data = [
                "\n".join(lines[i:i + chunk_size]) for i in range(0, len(lines), chunk_size)
            ]
            slides_data = slides_data[:slides]

        # Reconstruct text with proper slide headers for downstream logic
        reconstructed = ""
        for i, slide_text in enumerate(slides_data, 1):
            reconstructed += f"\nSlide {i}:\n{slide_text}\n"

        return reconstructed.strip()

    except Exception as e:
        raise RuntimeError(f"AI generation failed: {str(e)}")

    
# -----------------------------
# 🧩 PowerPoint Builder
# -----------------------------
def create_ppt(content, topic, color_hex, content_type="bullets"):
    from pptx.enum.text import PP_ALIGN
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]

    slides_data = re.split(r"(?:Slide\s*\d+[:\-]?)", content)
    for slide_text in slides_data:
        slide_text = slide_text.strip()
        if not slide_text or "Title:" not in slide_text:
            continue

        title_match = re.search(r"Title:\s*(.*)", slide_text)
        title = title_match.group(1).strip() if title_match else topic
        body_text = re.sub(r"Title:.*", "", slide_text, count=1).strip()

        if not body_text:
            continue

        slide = prs.slides.add_slide(slide_layout)
        apply_custom_color(slide, color_hex)

        # --- Title Box ---
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(8.5), Inches(1.3))
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        p = title_tf.add_paragraph()
        p.text = title
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.LEFT
        p.bullet = None 

        if len(title) <= 40:
            p.font.size = Pt(34)
        elif len(title) <= 70:
            p.font.size = Pt(28)
        else:
            p.font.size = Pt(24)

        # --- Content Box ---
        content_box = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(8.5), Inches(4.5))
        tf = content_box.text_frame
        tf.word_wrap = True

        lines = [l.strip() for l in body_text.split("\n") if l.strip()]

        for line in lines:
            # Limit paragraph size to ~4–5 lines (approx 220 chars)
            if len(line) > 120:
                line = line[:120] + "..."
            p = tf.add_paragraph()
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.space_after = Pt(6)
            p.alignment = PP_ALIGN.LEFT

            if content_type == "bullets":
                p.text = f"• {line}"
            else:  # paragraphs
                p.text = line
                p.bullet = None

    filename = f"{topic.replace(' ', '_')}_AI_Presentation.pptx"
    prs.save(filename)
    return filename


# -----------------------------
# 🌍 Glassy Modern UI + Fixed Logic
# -----------------------------
@app.route("/", methods=["GET", "POST"])
def index():
    html = """
    <!DOCTYPE html>
    <html lang="en">
    <head>
        <meta charset="UTF-8" />
        <meta name="viewport" content="width=device-width, initial-scale=1.0" />
        <title>AI PowerPoint Generator</title>
        <link href="https://fonts.googleapis.com/css2?family=Poppins:wght@300;400;600&display=swap" rel="stylesheet">

        <style>
            body {
                font-family: 'Poppins', sans-serif;
                background: radial-gradient(circle at 10% 20%, #0b132b, #1c2541);
                color: #fff;
                margin: 0;
                min-height: 100vh;
                display: flex;
                justify-content: center;
                align-items: center;
                overflow: hidden;
            }

            .circle {
                position: absolute;
                border-radius: 50%;
                background: linear-gradient(135deg, rgba(0,255,255,0.25), rgba(0,100,255,0.15));
                animation: float 10s ease-in-out infinite;
            }
            .circle:nth-child(1) { width: 320px; height: 320px; top: 15%; left: 10%; animation-delay: 0s; }
            .circle:nth-child(2) { width: 220px; height: 220px; bottom: 15%; right: 15%; animation-delay: 2.5s; }

            @keyframes float {
                0%, 100% { transform: translateY(0) rotate(0deg); }
                50% { transform: translateY(-25px) rotate(20deg); }
            }

            form {
                position: relative;
                background: rgba(255,255,255,0.06);
                border: 1px solid rgba(255,255,255,0.15);
                backdrop-filter: blur(20px);
                border-radius: 25px;
                padding: 45px 40px;
                box-shadow: 0 0 40px rgba(0,255,255,0.15);
                width: 90%;
                max-width: 480px;
                z-index: 2;
                animation: fadeInUp 0.8s ease;
            }

            @keyframes fadeInUp {
                from { opacity: 0; transform: translateY(40px); }
                to { opacity: 1; transform: translateY(0); }
            }

            h1 {
                text-align: center;
                font-weight: 600;
                margin-bottom: 30px;
                font-size: 28px;
                background: linear-gradient(90deg, #00f0ff, #00ffcc);
                -webkit-background-clip: text;
                -webkit-text-fill-color: transparent;
            }

            input, button {
                width: 100%;
                padding: 12px 15px;
                margin: 10px 0;
                border: none;
                border-radius: 12px;
                font-size: 15px;
                outline: none;
                transition: 0.3s;
            }

            input[type="text"], input[type="number"] {
                background: rgba(255,255,255,0.15);
                color: #fff;
            }

            input[type="text"]:focus, input[type="number"]:focus {
                background: rgba(255,255,255,0.25);
                box-shadow: 0 0 12px rgba(0,255,255,0.3);
            }

            input[type="color"] {
                height: 50px;
                border: 2px solid rgba(255,255,255,0.3);
                cursor: pointer;
                background: transparent;
            }

            .radio-group {
                display: flex;
                justify-content: space-between;
                color: #cceeff;
                margin: 15px 0;
            }

            button {
                background: linear-gradient(90deg, #00c6ff, #0072ff);
                color: white;
                font-weight: 600;
                letter-spacing: 0.3px;
                box-shadow: 0 0 20px rgba(0,255,255,0.3);
                cursor: pointer;
                transition: all 0.3s ease;
            }

            button:hover {
                transform: scale(1.05);
                box-shadow: 0 0 35px rgba(0,255,255,0.5);
            }

            #overlay, #success {
                display: none;
                position: fixed;
                top: 0; left: 0;
                width: 100%; height: 100%;
                background: rgba(0,0,0,0.65);
                backdrop-filter: blur(10px);
                z-index: 10;
                justify-content: center;
                align-items: center;
                flex-direction: column;
                color: #fff;
                text-align: center;
            }

            .loader {
                border: 6px solid rgba(255,255,255,0.2);
                border-top: 6px solid #00ffff;
                border-radius: 50%;
                width: 70px; height: 70px;
                animation: spin 1s linear infinite;
                margin-bottom: 20px;
            }

            @keyframes spin { 0% { transform: rotate(0deg); } 100% { transform: rotate(360deg); } }

            #success .box {
                background: rgba(255,255,255,0.1);
                padding: 40px 60px;
                border-radius: 20px;
                animation: fadeIn 0.8s ease;
            }

            @keyframes fadeIn { from {opacity:0;} to {opacity:1;} }
        </style>

       <script>
    function toggleFields(){
        let t=document.querySelector('input[name="content_type"]:checked').value;
        document.getElementById('bullet_count').style.display=(t==='bullets')?'block':'none';
        document.getElementById('para_count').style.display=(t==='paragraphs')?'block':'none';
    }

    function showLoading(){
        document.getElementById('overlay').style.display='flex';
        document.querySelector('#overlay p').textContent = "⚙️ Generating your PowerPoint... please wait ⏳";
        
        // After 6 seconds, check if download started successfully
        setTimeout(checkDownload, 6000);
    }

    async function checkDownload(){
        try {
            const res = await fetch('/done');
            if (res.ok) {
                showSuccess();
            }
        } catch (err) {
            console.log('Waiting for download to complete...');
            setTimeout(checkDownload, 2000);
        }
    }

    function showSuccess(){
        document.getElementById('overlay').style.display='none';
        document.getElementById('success').style.display='flex';
        document.querySelector("form").reset();
    }

    document.addEventListener('visibilitychange', () => {
        if (document.visibilityState === 'visible' && document.getElementById('overlay').style.display==='flex'){
            setTimeout(showSuccess, 1000);
        }
    });
</script>

    </head>

    <body>
        <div class="circle"></div>
        <div class="circle"></div>

        <div id="overlay">
            <div class="loader"></div>
            <p>⚙️ Generating your PowerPoint... please wait ⏳</p>
        </div>

        <div id="success">
            <div class="box">
                ✅ <strong>Your PowerPoint is ready!</strong><br><br>
                It has been downloaded successfully.<br>
                <button onclick="document.getElementById('success').style.display='none'">OK</button>
            </div>
        </div>

        <form method="POST" onsubmit="showLoading()">
            <h1>AI PowerPoint Generator</h1>

            <input type="text" name="topic" placeholder="Enter topic" required>
            <input type="number" name="slides" placeholder="Number of slides" min="1" required>

            <div class="radio-group">
                <label><input type="radio" name="content_type" value="bullets" checked onchange="toggleFields()"> Bullets</label>
                <label><input type="radio" name="content_type" value="paragraphs" onchange="toggleFields()"> Paragraphs</label>
            </div>

            <div id="bullet_count"><input type="number" name="count" placeholder="How many bullets?" min="1" max="8"></div>

            <label>Pick background color:</label>
            <input type="color" name="color" value="#0072ff">

            <button type="submit">Generate Presentation</button>
        </form>
    </body>
    </html>
    """

    if request.method == "POST":
        topic = request.form.get("topic")
        slides = int(request.form.get("slides"))
        color = request.form.get("color") or "#0072ff"
        content_type = request.form.get("content_type")

    # ✅ Enforce max 6
        user_count = int(request.form.get("count") or 4)
        count = min(user_count, 6)

        try:
            content = get_generated_content(topic, slides, "professional", content_type, count)
            pptx_file = create_ppt(content, topic, color, content_type)

            if not pptx_file or not os.path.exists(pptx_file):
                return "<h3 style='color:red;'>⚠️ PPT file could not be created. Please try again.</h3>"

            return send_file(pptx_file, as_attachment=True)

        except Exception as e:
            if "ConnectionError" in str(e) or "Timeout" in str(e):
                error_msg = "🌐 Network issue! Please check your internet and try again."
            elif "API" in str(e):
                error_msg = "🔑 API key error — check your .env configuration."
            else:
                error_msg = f"❌ Unexpected Error: {str(e)}"

            # Return same glass UI + error message box
            error_html = html + f"""
            <div style='position:fixed;bottom:20px;left:50%;transform:translateX(-50%);
                        background:rgba(255,50,50,0.2);padding:15px 30px;border-radius:15px;
                        color:#ff8080;font-weight:bold;backdrop-filter:blur(10px);
                        border:1px solid rgba(255,255,255,0.2);'>
                {error_msg}
            </div>"""
            return render_template_string(error_html)

    return render_template_string(html)


@app.route("/done")
def done():
    return "ok"


if __name__ == "__main__":
    app.run(debug=True)

